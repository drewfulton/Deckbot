#----------------------------------------------------------------------------
# Name:        presenters.py
# Purpose:     Application Logic
# Author:    Drew Fulton
# Created:    April 2020

import models, urllib.request, os, calendar, datetime, re, operator, statistics
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

class DeckbotPresenter(object):
	''' The main application logic.
	'''

	def __init__(self, view=None, company_id=None):
		models.get_token()
		self.view = view
		self.dir = os.path.dirname(os.path.abspath(__file__))
		if company_id is not None:
			company=models.Company(company_id)
		else:
			self.all_companies = models.get_all_companies(offline=True)
			company = self.init_view()
		self.get_company_details(company)
		self.create_deckbot(company)

	def init_view(self):
		''' Upon initial start, get the list of available companies, and prompt user.
		'''
		company = self.view.select_company(self.all_companies)
		return company

	def get_company_details(self, company):
		''' Get all the information and metrics for the company.
		'''
		company.get_company_overview()
		company.metrics = company.get_company_metrics()
		return company

	def create_deckbot(self, company):
		''' Genererate Powerpoint FactPack with 3 main slides.
		'''
		ppt = Presentation()

		ppt = self.build_title_slide(ppt, company)
		ppt = self.build_overview_slide(ppt, company)
		ppt = self.build_revenue_slide(ppt, company)   
	
		if not os.path.exists(f"{self.dir}/exports"):
			os.mkdir(f"{self.dir}/exports") 
				 
		ppt.save(f"{self.dir}/exports/{company.name}.pptx")
	
		print(f"Please find your file at /exports/{company.name}.pptx")
	
	
	def build_title_slide(self, ppt, company):
		''' Builds the title slide using existing title/subtitle placeholders.
		Adds the company logo and photo to the slide.
		'''
		# Title Slide - Title Slide Layout
		title_slide_layout = ppt.slide_layouts[0]
		title_slide = ppt.slides.add_slide(title_slide_layout)
		title = title_slide.shapes.title
		subtitle = title_slide.placeholders[1]
		today = datetime.date.today().strftime("%B %d, %Y")
	
		logofile, extension = os.path.splitext(company.logoUrl)
		try:
			urllib.request.urlretrieve(
				company.logoUrl, 
				f"{self.dir}/assets/{company.name}{extension}"
			)
			logo = title_slide.shapes.add_picture(
				f"{self.dir}/assets/{company.name}{extension}", 
				Inches(.5), 
				Inches(.5), 
				height=Inches(1.5)
			)
			os.remove(f"{self.dir}/assets/{company.name}{extension}")
		except:
			print("Error retrieving logo from website and will be skipped.")
	
		latestRev = f"Data through Q{company.latestRevenue['quarter']} {company.latestRevenue['year']}"
		title.text = f"{company.name} - Factpack"
		subtitle_frame = subtitle.text_frame
		subtitle_frame.clear()
	
		p1 = subtitle_frame.paragraphs[0]
		p1.text = f"Created on {today}"
		p2 = subtitle_frame.add_paragraph()
		p2.text = latestRev
	
		img = title_slide.shapes.add_picture(
			f"{self.dir}/assets/cover.jpg", 
			Inches(0), 
			Inches(5.5), 
			height=Inches(2)
		)
	
		return ppt

	def build_overview_slide(self, ppt, company):
		''' Creates an Overview slide from a Blank slide template.  
		General Description is broken into paragraphs for each sentence and placed
		in one column.  Basic facts are placed in a second text box.
		'''
		# Overview Slide - Blank Slide Layout
		overview_slide_layout = ppt.slide_layouts[6]
		overview_slide = ppt.slides.add_slide(overview_slide_layout)
		shapes = overview_slide.shapes
	
		# Create Title Box
		title_sizes = {}
		title_sizes["left"] = Inches(0)
		title_sizes["top"] = Inches(0.5)
		title_sizes["width"] = Inches(10)
		title_sizes["height"] = Inches (1)
		title_text = [f"{company.name} Introduction"]
		shapes = self.create_textbox(
			shapes, 
			title_sizes, 
			title_text, 
			alignment=PP_ALIGN.CENTER, 
			color=RGBColor(0,0,0), 
			size=Pt(36)
		)
	
		# Create Summary Title Box
		sum_title_sizes = {}
		sum_title_sizes["left"] = Inches(.5)
		sum_title_sizes["top"] = Inches(1.5)
		sum_title_sizes["width"] = Inches(4.5)
		sum_title_sizes["height"] = Inches (1)
		sum_title_text = "Company Summary"
		shapes = self.create_textbox(
			shapes, 
			sum_title_sizes, 
			sum_title_text,
			alignment=PP_ALIGN.CENTER, 
			color=RGBColor(0,0,0), 
			size=Pt(24)
		)

		# Create Summary Box
		summary_sizes = {}
		summary_sizes["left"] = Inches(.5)
		summary_sizes["top"] = Inches(2)
		summary_sizes["width"] = Inches(4.5)
		summary_sizes["height"] = Inches (5)    
		summary_text = split_into_sentences(company.description)
		shapes = self.create_textbox(
			shapes, 
			summary_sizes, 
			summary_text, 
			alignment=PP_ALIGN.LEFT, 
			size=Pt(10)
		)

		# Create Facts Title Box
		fact_title_sizes = {}
		fact_title_sizes["left"] = Inches(5)
		fact_title_sizes["top"] = Inches(1.5)
		fact_title_sizes["width"] = Inches(4.5)
		fact_title_sizes["height"] = Inches (1)
		fact_title_text = "Company Facts"
		shapes = self.create_textbox(
			shapes, 
			fact_title_sizes, 
			fact_title_text, 
			alignment=PP_ALIGN.CENTER, 
			color=RGBColor(0,0,0), 
			size=Pt(24)
		)

		# Create Facts Box
		facts_sizes = {}
		facts_sizes["left"] = Inches(5)
		facts_sizes["top"] = Inches(2)
		facts_sizes["width"] = Inches(4.5)
		facts_sizes["height"] = Inches (5)
		facts_text = []
		facts_text.append((
			("Revenue: ", {"bold":True}), 
			(f"US$ {company.latestRevenue['valueUSD']} bn", {"bold":False})
		))
		facts_text.append((
			("Employees: ", {"bold":True}), 
			(f"{'{:,d}'.format(int(company.employees*1000000))}", {"bold":False})
		))
		facts_text.append((
			("Currency: ", {"bold":True}), 
			(f"{company.currency}", {"bold":False})
		))
		facts_text.append((
			("Type: ", {"bold":True}), 
			(f"{company.type}", {"bold":False})
		))
		facts_text.append((
			("Website: ", {"bold":True}), 
			(f"{company.website}", {"bold":False, "url":company.website})
		))
		facts_text.append((
			("Headquarters: ", {"bold":True}), 
			(f"{company.address}", {"bold":False})
		))
		facts_text.append((
			("Current Quarter: ", {"bold":True}), 
			(f"{company.currentQuarter}", {"bold":False})
		))
		facts_text.append((
			("Current Quarter Ends: ", {"bold":True}), 
			(company.quarterEnd[0:10] , {"bold":False})
		))
		facts_text.append((
			("Fiscal Year End: ", {"bold":True}), 
			(calendar.month_name[company.fiscalYearEnd], {"bold":False})
		))
		shapes = self.create_textbox(
			shapes, 
			facts_sizes, 
			facts_text, 
			alignment=PP_ALIGN.LEFT, 
			size=Pt(10), 
			space_after=Pt(10)
		)

		return ppt

	def build_revenue_slide(self, ppt, company):
		''' Creates a Revenue Metrics slide with two charts and corresponding info.
		First chart compares the company's revenue from teh last 3 years against the 
		sector's median revenues.  Second chart compares the company's Latest Revenue
		against its peers.  Basic analysis for each chart is included.
		'''
	
		metrics_slide_layout = ppt.slide_layouts[6]
		metrics_slide = ppt.slides.add_slide(metrics_slide_layout)
		shapes = metrics_slide.shapes
	
		# Get Metric Details
		for m in company.metrics:
			if m.name == "Revenue":
				revenue = m
				break
	
		revenue.get_metric_details()
	
		# Title Box
		title_sizes = {}
		title_sizes["left"] = Inches(0)
		title_sizes["top"] = Inches(0.25)
		title_sizes["width"] = Inches(10)
		title_sizes["height"] = Inches (.5)
		title_text = [f"{company.name} {revenue.name} Metrics"]
		shapes = self.create_textbox(
			shapes, 
			title_sizes, 
			title_text, 
			alignment=PP_ALIGN.CENTER, 
			color=RGBColor(0,0,0), 
			size=Pt(36)
		)

		# Summary Box
		sum_sizes = {}
		sum_sizes["left"] = Inches(.5)
		sum_sizes["top"] = Inches(1)
		sum_sizes["width"] = Inches(9)
		sum_sizes["height"] = Inches (1)
		sum_text = revenue.description.replace('\n', ' ')
		shapes = self.create_textbox(
			shapes, 
			sum_sizes, 
			sum_text, 
			alignment=PP_ALIGN.CENTER, 
			color=RGBColor(0,0,0), 
			size=Pt(10)
		)
	
	
		# Latests Revenue vs Peers Section
		latest_rev_v_peers_chart_data = CategoryChartData()
		latest_rev_v_peers_chart_data.categories = ['Latest Values']
		data = {}
		categories = []
		values = []
		i=1
		for c in revenue.chart[0]["companies"]:
			data[c["name"]] = c["data"][0]["value"]
			if i==1:
				period_label = c["data"][0]["label"]
				period = c["data"][0]["period"]
			i=+1
		sorted_data = sorted(data.items(), key=operator.itemgetter(1))    
		for d in sorted_data:
			categories.append(d[0])
			values.append(d[1])
			if d[0] == company.name:
				company_position = sorted_data.index(d)+1
		
		latest_rev_v_peers_chart_data.categories = categories
		latest_rev_v_peers_chart_data.add_series(
			"Latest Values", 
			values, 
			number_format="#,###.#"
		)
		x, y, cx, cy = Inches(5.25), Inches(2.5), Inches(4.5), Inches(4)
		graphic_frame = metrics_slide.shapes.add_chart(    
			XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, latest_rev_v_peers_chart_data
		)
		chart = graphic_frame.chart    
		plot = chart.plots[0]
		plot.has_data_labels = True
		data_labels = plot.data_labels
		data_labels.font.size = Pt(8)
		data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
		cat_axis = chart.category_axis
		cat_axis.tick_labels.font.size = Pt(8)
		val_axis = chart.value_axis
		val_axis.tick_labels.font.size = Pt(12)
		val_axis.visible = False

		# Add Labels & Textboxes
		latest_rev_title_text = "Latest revenues vs. peers"
		latest_rev_title_sizes = {}
		latest_rev_title_sizes["left"] = Inches(5.25)
		latest_rev_title_sizes["top"] = Inches(1.5)
		latest_rev_title_sizes["width"] = Inches(4.5)
		latest_rev_title_sizes["height"] = Inches (.5)
		shapes = self.create_textbox(
			shapes, 
			latest_rev_title_sizes, 
			latest_rev_title_text, 
			alignment=PP_ALIGN.LEFT, 
			color=RGBColor(0,0,0), 
			size=Pt(24)
		)

		latest_rev_det_text = []
		latest_rev_det_text.append((("Revenue (USD bn)", {"bold":True, "size": Pt(14)}),))
		latest_rev_det_text.append(((f"{period_label}", {"bold":False, "size": Pt(12)}),))
		latest_rev_det_sizes = {}
		latest_rev_det_sizes["left"] = Inches(5.25)
		latest_rev_det_sizes["top"] = Inches(2)
		latest_rev_det_sizes["width"] = Inches(4.5)
		latest_rev_det_sizes["height"] = Inches (.5)
		shapes = self.create_textbox(
			shapes, 
			latest_rev_det_sizes, 
			latest_rev_det_text, 
			alignment=PP_ALIGN.LEFT, 
			color=RGBColor(0,0,0), 
			size=Pt(14)
		)
	
		# Analysis Text
		period_label_plain = {}
		period_label_plain["ltm"] = "last twelve months"
		period_label_plain["y"] = "year"
		company_count = len(sorted_data)
		quartile = get_quartile(company_position/company_count)
		latest_rev_anal_text = f"{company.name} Revenue is in the {quartile} quartile of the peer group over the {period_label_plain[period]}."
		latest_rev_anal_sizes = {}
		latest_rev_anal_sizes["left"] = Inches(5.25)
		latest_rev_anal_sizes["top"] = Inches(6.5)
		latest_rev_anal_sizes["width"] = Inches(4.5)
		latest_rev_anal_sizes["height"] = Inches (.5)
		shapes = self.create_textbox(
			shapes, 
			latest_rev_anal_sizes, 
			latest_rev_anal_text, 
			alignment=PP_ALIGN.LEFT, 
			color=RGBColor(0,0,0), 
			size=Pt(14)
		)


		# Revenue for Last 3 Years Section
		last_three_chart_data = CategoryChartData()
		last_three_data = {}
		for c in revenue.chart[0]["companies"]:
			last_three_data[c["name"]] = {}
			for p in c["data"]:
				if "Last 3 years" in p["groups"]:
					last_three_data[c["name"]][p["label"]] = p["value"]
			
		cats = list(last_three_data[company.name].keys())
		cats.reverse()
		last_three_chart_data.categories = cats
		company_data = []
		median_data = {}
		medians = {}
		for cat in cats:
			company_data.append(last_three_data[company.name][cat])
			median_data[cat] = [] 
			for comp in last_three_data:
				median_data[cat].append(last_three_data[comp][cat])
			medians[cat] = statistics.median(median_data[cat])
	
		med = []
		for m in medians:
			med.append(medians[m])
		last_three_chart_data.add_series(
			company.name, 
			company_data, 
			number_format="#,###.#"
		)
		last_three_chart_data.add_series("Medians", med, number_format="#,###.#")
		x, y, cx, cy = Inches(.5), Inches(2.5), Inches(4.5), Inches(4)
		last_three_frame = metrics_slide.shapes.add_chart( 
			XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, last_three_chart_data
		)
	
		chart = last_three_frame.chart    
		plot = chart.plots[0]
		plot.has_data_labels = True
		data_labels = plot.data_labels
		data_labels.font.size = Pt(8)
		data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
		cat_axis = chart.category_axis
		cat_axis.tick_labels.font.size = Pt(8)
		val_axis = chart.value_axis
		val_axis.tick_labels.font.size = Pt(12)
		val_axis.visible = False
		chart.has_legend = True
		chart.legend.position = XL_LEGEND_POSITION.BOTTOM
		chart.legend.include_in_layout = False
		chart.legend.font.size = Pt(12)

		# Add Labels & Textboxes
		last_three_rev_title_text = "Revenues (last 3 years)"
		last_three_rev_title_sizes = {}
		last_three_rev_title_sizes["left"] = Inches(0.5)
		last_three_rev_title_sizes["top"] = Inches(1.5)
		last_three_rev_title_sizes["width"] = Inches(4.5)
		last_three_rev_title_sizes["height"] = Inches (.5)
		shapes = self.create_textbox(
			shapes, 
			last_three_rev_title_sizes, 
			last_three_rev_title_text, 
			alignment=PP_ALIGN.LEFT, 
			color=RGBColor(0,0,0), 
			size=Pt(24)
		)

		last_three_rev_det_text = []
		last_three_rev_det_text.append((
			("Revenue (USD bn)", {"bold":True, "size": Pt(14)}),
		))
		last_three_rev_det_text.append((
			(f"{cats[0]} - {cats[-1]}", {"bold":False, "size": Pt(12)}),
		))
		last_three_rev_det_sizes = {}
		last_three_rev_det_sizes["left"] = Inches(0.5)
		last_three_rev_det_sizes["top"] = Inches(2)
		last_three_rev_det_sizes["width"] = Inches(4.5)
		last_three_rev_det_sizes["height"] = Inches (.5)
		shapes = self.create_textbox(
			shapes, 
			last_three_rev_det_sizes, 
			last_three_rev_det_text, 
			alignment=PP_ALIGN.LEFT, 
			color=RGBColor(0,0,0), 
			size=Pt(14)
		)
	
		# Analysis Text
		period_label_plain = {}
		period_label_plain["ltm"] = "last twelve months"
		period_label_plain["y"] = "year"
	
		if company_data[-1] > company_data[0]:
			direction = "grew"
		elif company_data[0] == company_data[-1]:
			direction = "held steady"
		else:
			direction = "fell"

		company_count = len(sorted_data)
		quartile = get_quartile(company_position/company_count)
	
		last_three_rev_anal_text = f"{company.name} Revenue {direction} from USD{company_data[0]}B in FY{cats[0]} to USD{company_data[-1]}B at the end of Q{company.latestRevenueGrowth['quarter']} FY{company.latestRevenueGrowth['year']}."
		last_three_rev_anal_sizes = {}
		last_three_rev_anal_sizes["left"] = Inches(.5)
		last_three_rev_anal_sizes["top"] = Inches(6.5)
		last_three_rev_anal_sizes["width"] = Inches(4.5)
		last_three_rev_anal_sizes["height"] = Inches (.5)
		shapes = self.create_textbox(
			shapes, 
			last_three_rev_anal_sizes, 
			last_three_rev_anal_text, 
			alignment=PP_ALIGN.LEFT, 
			color=RGBColor(0,0,0), 
			size=Pt(14)
		)

		return ppt

	def set_shape_colors(
		self, 
		shape, 
		line_color=RGBColor(255,255,255), 
		fill_color=RGBColor(255,255,255), 
		line_width = Pt(0)
		):
		''' Set colors for fill and lines for a given shape
		'''
		fill = shape.fill
		fill.solid()
		fill.fore_color.rgb = fill_color
		line = shape.line
		line.color.rgb = line_color
		line.width = line_width
	
		return shape


	def create_textbox(
		self, 
		shapes, 
		dim, 
		content, 
		space_after=None, 
		color=RGBColor(0,0,0), 
		wordwrap=True, 
		alignment=PP_ALIGN.CENTER, 
		auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
		size=None
		):
		''' Creates a text box given the dimensions, content, and formatting
		'''        
		shape = shapes.add_textbox(dim["left"], dim["top"], dim["width"], dim["height"])
		text_frame = shape.text_frame
		text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE    
		text_frame.clear()
		text_frame.word_wrap = wordwrap
		i = 0
		if isinstance(content, str):
			# Single String paragraph
			p = text_frame.paragraphs[0]
			p.alignment = alignment
			p.space_after=space_after
			run = p.add_run()
			run.text = content
			font = run.font
			font.color.rgb = color
			font.size = size
		else:
			# Multiple paragraph
			for para in content:
				if i == 0:
					p = text_frame.paragraphs[0]
				else:
					p = text_frame.add_paragraph()
				i =+ 1
				p.alignment = alignment
				p.space_after = space_after
				if type(para) is not str:
					# Multiple character formating in a single paragraph
					for r in para:
						run = p.add_run()
						run.text = r[0]
						font = run.font
						if "bold" in r[1]:
							font.bold = r[1]["bold"]
						if "italic" in r[1]:
							font.italic = r[1]["italic"]
						if "color" in r[1]:
							font.color.rgb = r[1]["color"]
						if "size" in r[1]:
							font.size = r[1]["size"]
						if "url" in r[1]:
							run.hyperlink.address = r[1]["url"]            
				else: 
					# Just a plain formatted paragraph
					run = p.add_run()
					run.text = para
					font = run.font
					font.color.rgb = color
					font.size = size
	
		return shapes

def get_quartile(percentage):
	''' Calculate quartile given a percentage. '''
	if percentage > .75:
		quartile = "top"
	elif .5 > percentage >=.75:
		quartile = "second"
	elif .25 > percentage >= .5:
		quartile = "third"
	else:
		quartile = "last"
	return quartile


# Split Paragraphs into Sentences 
alphabets= "([A-Za-z])"
prefixes = "(Mr|St|Mrs|Ms|Dr)[.]"
suffixes = "(Inc|Ltd|Jr|Sr|Co)"
starters = "(Mr|Mrs|Ms|Dr|He\s|She\s|It\s|They\s|Their\s|Our\s|We\s|But\s|However\s|That\s|This\s|Wherever)"
acronyms = "([A-Z][.][A-Z][.](?:[A-Z][.])?)"
websites = "[.](com|net|org|io|gov)"

def split_into_sentences(text):
	''' Splits a paragraph into a list of sentences.'''
	text = " " + text + "  "
	text = text.replace("\n"," ")
	text = re.sub(prefixes,"\\1<prd>",text)
	text = re.sub(websites,"<prd>\\1",text)
	if "Ph.D" in text: text = text.replace("Ph.D.","Ph<prd>D<prd>")
	text = re.sub("\s" + alphabets + "[.] "," \\1<prd> ",text)
	text = re.sub(acronyms+" "+starters,"\\1<stop> \\2",text)
	text = re.sub(alphabets + "[.]" + alphabets + "[.]" + alphabets + "[.]","\\1<prd>\\2<prd>\\3<prd>",text)
	text = re.sub(alphabets + "[.]" + alphabets + "[.]","\\1<prd>\\2<prd>",text)
	text = re.sub(" "+suffixes+"[.] "+starters," \\1<stop> \\2",text)
	text = re.sub(" "+suffixes+"[.]"," \\1<prd>",text)
	text = re.sub(" " + alphabets + "[.]"," \\1<prd>",text)
	if "\"" in text: text = text.replace(".\"","\".")
	if "!" in text: text = text.replace("!\"","\"!")
	if "?" in text: text = text.replace("?\"","\"?")
	text = text.replace(".",".<stop>")
	text = text.replace("?","?<stop>")
	text = text.replace("!","!<stop>")
	text = text.replace("<prd>",".")
	sentences = text.split("<stop>")
	sentences = sentences[:-1]
	sentences = [s.strip() for s in sentences]
	return sentences